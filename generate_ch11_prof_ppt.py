from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def generate_prof_ppt(output_path, template_path, img_dir, text_file):
    # Load template or create new
    if os.path.exists(template_path):
        prs = Presentation(template_path)
        # Clear template slides
        for _ in range(len(prs.slides)):
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
    else:
        prs = Presentation()

    def add_prof_slide(title, points, img=None, img_scale=6):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        
        # Consistent Header
        header_text = "十一五规划教材                                          《信息系统项目管理》讲义"
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = header_text
        p.font.size = Pt(12)
        slide.shapes.add_connector(1, Inches(0.5), Inches(0.7), Inches(9.5), Inches(0.7))

        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        # Add points
        for i, pt in enumerate(points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = pt
            # Detect nested points (starting with - or whitespace)
            if pt.startswith("  ") or pt.startswith("- "):
                p.level = 1
                p.text = pt.lstrip(" -")
            else:
                p.level = 0
            p.font.size = Pt(18)

        if img:
            img_path = os.path.join(img_dir, img)
            if os.path.exists(img_path):
                # If there are few points, place image in middle. If many, bottom.
                top = Inches(2.5) if len(points) < 3 else Inches(4.5)
                slide.shapes.add_picture(img_path, Inches(1), top, width=Inches(img_scale))

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tx = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = "信息系统项目管理\n第11章 项目成本管理"
    p.font.bold = True
    p.font.size = Pt(44)
    p.alignment = PP_ALIGN.CENTER

    # Content - Based on Chapter Outline
    # 11.1
    add_prof_slide("11.1 项目成本管理概述", [
        "项目成本管理是确定如何估算、预算、管理、监督和控制项目成本的过程。",
        "主要作用：在整个项目期间为如何管理项目成本提供指南和方向。",
        "核心工作：",
        "  - 规划成本管理：确定管理基准、框架和流程。",
        "  - 估算成本：近似估算完成项目工作所需资源成本。",
        "  - 制定预算：汇总所有单个活动或工作包的估算成本，建立成本基准。",
        "  - 控制成本：监督项目状态，更新成本和管理成本基准变更。"
    ])

    # 11.3 Plan Cost Management
    add_prof_slide("11.3 规划成本管理", [
        "规划成本管理是确定如何估算、预算、管理、监督和控制项目成本的过程。",
        "应该在项目规划阶段的早期就对成本管理工作进行规划。",
        "主要作用：为整个项目期间如何管理项目成本提供指南和方向。"
    ], img="fig_11_1.png", img_scale=7)

    add_prof_slide("11.3 规划成本管理 - ITTO", [
        "规划成本管理过程汇总表："
    ], img="itto_11_3_v2.png", img_scale=8)

    add_prof_slide("11.3.1 输入", [
        "1. 项目章程：规定了预先批准的财务资源及审批要求。",
        "2. 项目管理计划：",
        "  - 进度管理计划：提供了影响成本估算和管理的过程及控制方法。",
        "  - 风险管理计划：提供了影响成本估算和管理的过程及控制方法。",
        "3. 事业环境因素：市场条件、货币汇率、发布的商业信息等。",
        "4. 组织过程资产：财务控制程序、历史信息知识库、财务数据库等。"
    ])

    add_prof_slide("11.3.3 成本管理计划内容", [
        "成本管理计划中一般需要规定：",
        "  - 计量单位：如人时数、米、吨、立方码、总价等。",
        "  - 精确度：设定成本估算向上或向下取整的程度。",
        "  - 准确度：为活动成本估算规定一个可接受的区间（如±10%）。",
        "  - 组织程序链接：WBS为成本管理计划提供了框架（控制账户 CA）。",
        "  - 控制临界值：偏差临界值临界值，用于监督成本绩效。",
        "  - 绩效测量规则：规定用于绩效测量的挣值管理（EVM）规则。"
    ])

    # 11.4 Estimate Costs
    add_prof_slide("11.4 估算成本", [
        "估算成本是对完成项目工作所需资源成本进行近似估算的过程。",
        "本过程的主要作用是确定项目所需的资金。",
        "成本预测是在某特定时点根据已知信息所做出的的成本预测。"
    ], img="fig_11_2.png", img_scale=7)

    add_prof_slide("11.4 估算成本 - ITTO", [
        "估算成本过程汇总表："
    ], img="itto_11_4_v2.png", img_scale=8)

    add_prof_slide("11.4 估算准确度级别", [
        "随着项目进展，项目估算的准确性将逐步提高：",
        "  - 启动阶段：粗略量级估算 (ROM)，区间为 -25% 到 +75%。",
        "  - 后期/确定性估算：区间可缩小至 -5% 到 +10%。",
        "考虑资源：人工、材料、设备、服务、设施，以及通货膨胀补贴、融资成本等。"
    ])

    add_prof_slide("11.4.2 工具与技术", [
        "1. 类比估算：使用以往类似项目的参数值或属性来估算。",
        "2. 参数估算：利用历史数据之间的统计关系和其它变量进行估算。",
        "3. 自下而上估算：对工作组成部分进行最精细估算并向上汇总。",
        "4. 三点估算：(乐观+4*可能+悲观)/6，考虑不确定性和风险。",
        "5. 储备分析：对应急储备（已知-未知风险）进行分析。"
    ])

    # 11.5 Determine Budget
    add_prof_slide("11.5 制定预算", [
        "制定预算是汇总所有单个活动或工作包的估算成本，建立一个经批准的成本基准的过程。",
        "主要作用：确定成本基准，用以衡量和监控项目绩效。"
    ], img="fig_11_3.png", img_scale=7)

    add_prof_slide("11.5 制定预算 - ITTO", [
        "制定预算过程汇总表："
    ], img="itto_11_5_v2.png", img_scale=8)

    add_prof_slide("11.5.2 制定预算：工具与技术", [
        "1. 成本汇总：从下到上逐层汇总到工作包和控制账户。",
        "2. 专家判断：根据历史类似项目进行判断。",
        "3. 历史审核：检查项目间成本关系及其它特征。",
        "4. 资金限额平衡：根据资金限制对支出进行平衡和调整。",
        "5. 融资：为项目获取外部资金。"
    ])

    add_prof_slide("11.5.3 项目预算的组成", [
        "项目预算 = 成本基准 + 管理储备。",
        "成本基准 = 活动估算 + 应急储备。",
        "管理储备用于“未知-未知”风险，不属于成本基准。"
    ], img="fig_11_4.png", img_scale=7)

    # 11.6 Control Costs
    add_prof_slide("11.6 控制成本", [
        "控制成本是监督项目状态，以更新项目成本和管理成本基准变更的过程。",
        "主要作用：发现实际支出与计划的偏差，并采取纠正措施。"
    ], img="fig_11_6.png", img_scale=7)

    add_prof_slide("11.6 控制成本 - ITTO", [
        "控制成本过程汇总表："
    ], img="itto_11_6_v2.png", img_scale=8)

    # EVM Sections (The "Killer" Professional Section)
    add_prof_slide("11.6.2 数据分析：挣值分析 (EVM)", [
        "挣值分析将范围、进度和资源绩效综合起来，衡量项目绩效和进度。",
        "核心指标：",
        "  - PV (Planned Value)：计划完成工作的预算。",
        "  - EV (Earned Value)：实际完成工作的预算。",
        "  - AC (Actual Cost)：实际发生的成本。"
    ], img="fig_11_7.png", img_scale=7)

    add_prof_slide("挣值分析（EVM）评价指标汇总 (1)", [
        "详细公示与定义 (PV, EV, AC, BAC, CV, SV, VAC)："
    ], img="evm_table_p1.png", img_scale=9)

    add_prof_slide("挣值分析（EVM）评价指标汇总 (2)", [
        "效率指数与预测指标 (CPI, SPI, EAC, ETC, TCPI)："
    ], img="evm_table_p2.png", img_scale=9)

    add_prof_slide("11.6.2 预测分析 (EAC)", [
        "完工估算 (EAC) 的不同计算情况：",
        "1. 典型偏差（之后按预算干）：EAC = AC + (BAC - EV)",
        "2. 非典型偏差（按当前效率干）：EAC = BAC / CPI",
        "3. 考虑进度影响：EAC = AC + [(BAC - EV) / (CPI * SPI)]",
        "VAC (完工偏差) = BAC - EAC。"
    ])

    prs.save(output_path)
    print(f"Professional PPT saved to {output_path}")

if __name__ == "__main__":
    out_file = r"e:\教学管理\第11章 项目成本管理_完整专业版.pptx"
    tmpl_file = r"e:\教学管理\第4章 信息系统项目的进度管理.pptx"
    img_dir = r"e:\教学管理\ch11_images"
    text_file = r"e:\教学管理\ch11_prof_text_final.txt"
    generate_prof_ppt(out_file, tmpl_file, img_dir, text_file)
