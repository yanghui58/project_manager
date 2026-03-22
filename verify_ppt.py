from pptx import Presentation
import os

def verify_ppt(path):
    if not os.path.exists(path):
        print(f"File not found: {path}")
        return
    
    prs = Presentation(path)
    print(f"Total slides: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides):
        title = slide.shapes.title.text if hasattr(slide, "shapes") and slide.shapes.title else "No Title"
        print(f"Slide {i+1}: {title}")

if __name__ == "__main__":
    verify_ppt(r"e:\教学管理\第11章 项目成本管理.pptx")
