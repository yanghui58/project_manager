from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def create_pure_textbook_ppt(output_path, template_path, img_dir, text_file):
    if os.path.exists(template_path):
        prs = Presentation(template_path)
        for _ in range(len(prs.slides)):
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
    else:
        prs = Presentation()

    def add_slide(title, paragraphs, img=None, img_scale=7.5):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        
        header_text = "十一五规划教材                                          《信息系统项目管理》讲义"
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = header_text
        p.font.size = Pt(12)
        slide.shapes.add_connector(1, Inches(0.5), Inches(0.7), Inches(9.5), Inches(0.7))

        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        for i, text in enumerate(paragraphs):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = text
            if text.startswith("  ") or text.startswith("●") or text.startswith("·"):
                p.level = 1
                p.text = text.lstrip(" ●·")
            else:
                p.level = 0
            p.font.size = Pt(17)

        if img:
            img_path = os.path.join(img_dir, img)
            if os.path.exists(img_path):
                top_margin = Inches(4.0) if len(paragraphs) > 3 else Inches(2.0)
                slide.shapes.add_picture(img_path, Inches(1), top_margin, width=Inches(img_scale))

    # Parse OCR Text File
    sections = {}
    current_section = "Intro"
    sections[current_section] = []
    
    with open(text_file, 'r', encoding='utf-8') as f:
        # Simple parsing logic looking for chapter headers
        for line in f:
            line = line.strip()
            if not line: continue
            if line.startswith("--- PDF Page"): continue
            if "信息系统项目管理师教程" in line: continue
            if line.startswith(("11.1", "11.2", "11.3", "11.4", "11.5", "11.6")):
                # Normalize section headers slightly for grouping
                sec_prefix = line[:4] 
                if sec_prefix in ["11.1", "11.2", "11.3", "11.4", "11.5", "11.6"]:
                    current_section = sec_prefix
                    if current_section not in sections:
                        sections[current_section] = []
            
            sections[current_section].append(line)

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tx = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = "信息系统项目管理师教程（第4版）\n第11章 项目成本管理\n【纯原版教材版】"
    p.font.bold = True
    p.font.size = Pt(44)
    p.alignment = PP_ALIGN.CENTER
    
    # Generate Slides from Content
    for sec_prefix, content in sections.items():
        if sec_prefix == "Intro": continue # Skip garbage at start
        
        # Batch content into slides (e.g. 7-8 lines per slide)
        chunk_size = 7
        
        # Add ITTO Tables explicitly based on section
        if sec_prefix == "11.3":
            add_slide("11.3 规划成本管理 - ITTO", ["规划成本管理过程汇总表"], img="itto_11_3_v2.png", img_scale=8.5)
        elif sec_prefix == "11.4":
            add_slide("11.4 估算成本 - ITTO", ["估算成本过程汇总表"], img="itto_11_4_v2.png", img_scale=8.5)
        elif sec_prefix == "11.5":
            add_slide("11.5 制定预算 - ITTO", ["制定预算过程汇总表"], img="itto_11_5_v2.png", img_scale=8.5)
        elif sec_prefix == "11.6":
            add_slide("11.6 控制成本 - ITTO", ["控制成本过程汇总表"], img="itto_11_6_v2.png", img_scale=8.5)

        for i in range(0, len(content), chunk_size):
            chunk = content[i:i+chunk_size]
            title = f"{sec_prefix} 内容详情 ({(i//chunk_size)+1})"
            
            # Map diagrams
            img = None
            if sec_prefix == "11.3" and i == 0: img = "fig_11_1.png"
            if sec_prefix == "11.4" and i == 0: img = "fig_11_2.png"
            if sec_prefix == "11.5" and i == 0: img = "fig_11_3.png"
            if sec_prefix == "11.5" and i == chunk_size*4: img = "fig_11_4.png"
            if sec_prefix == "11.6" and i == 0: img = "fig_11_6.png"

            # Filter out empty chunks
            if any(len(x.strip()) > 3 for x in chunk):
                add_slide(title, chunk, img=img)

            # Insert EVM Tables towards the end of 11.6
            if sec_prefix == "11.6" and i == chunk_size*6:
                 add_slide("11.6 挣值、计划价值和实际成本", [], img="fig_11_7.png")
            if sec_prefix == "11.6" and i == chunk_size*8:
                 add_slide("11.6 表11-2挣值计算汇总表 (1)", [], img="evm_table_p1.png", img_scale=9)
                 add_slide("11.6 表11-2挣值计算汇总表 (2)", [], img="evm_table_p2.png", img_scale=9)

    prs.save(output_path)
    print(f"Pure Textbook PPT saved to {output_path}")

if __name__ == "__main__":
    out_file = r"e:\教学管理\第11章 项目成本管理_纯教材版.pptx"
    tmpl_file = r"e:\教学管理\第4章 信息系统项目的进度管理.pptx"
    img_dir = r"e:\教学管理\ch11_images"
    text_file = r"e:\教学管理\ch11_ocr_text.txt"
    create_pure_textbook_ppt(out_file, tmpl_file, img_dir, text_file)
