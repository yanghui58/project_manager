from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_ch11_ppt(output_path, template_path):
    # Load from template to preserve master slides and styles
    if os.path.exists(template_path):
        prs = Presentation(template_path)
        # Clear existing slides
        for _ in range(len(prs.slides)):
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
    else:
        prs = Presentation()

    def add_standard_slide(title_text, content_items=None, subtitle=None):
        # Using "标题和内容" layout (index 1 usually)
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set Title (Rectangle 2 in user's PPT)
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = title_text
            
        # Set Header (Rectangle 4/5 in user's PPT)
        # We manually add it if not in master
        header_text = "十一五规划教材                                          《信息系统项目管理》讲义"
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = header_text
        p.font.size = Pt(12)
        
        # Add Horizontal Line
        slide.shapes.add_connector(1, Inches(0.5), Inches(0.7), Inches(9.5), Inches(0.7))

        # Set Content (Rectangle 3 in user's PPT)
        if content_items:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.text = content_items[0]
            for item in content_items[1:]:
                p = tf.add_paragraph()
                p.text = item
                p.level = 0
        
        if subtitle:
            # Similar to how the user puts process name at the top of content
            pass

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    tx1 = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
    tf1 = tx1.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "信息系统项目管理"
    p1.font.bold = True
    p1.font.size = Pt(44)
    p1.alignment = PP_ALIGN.CENTER

    tx2 = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(6), Inches(2))
    tf2 = tx2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "第11章 项目成本管理"
    p2.font.size = Pt(32)
    p2.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Overview
    add_standard_slide("项目成本管理概述", [
        "项目成本管理是为了保证项目在批准的预算内完成而对成本进行规划、估算、预算、融资、筹资、管理和控制的过程。",
        "核心过程包括：",
        "1. 规划成本管理：制定成本管理计划。",
        "2. 估算成本：对完成项目活动所需支出进行近似估算。",
        "3. 制定预算：汇总所有单个活动或工作包的估算成本，建立成本基准。",
        "4. 控制成本：监督项目状态，更新项目成本，管理成本基准变更。"
    ])

    # --- 11.3 规划成本管理 ---
    add_standard_slide("11.3 规划成本管理", [
        "规划成本管理是确定如何估算、预算、管理、监督和控制项目成本的过程。",
        "主要作用：在整个项目期间为如何管理项目成本提供指南和方向。"
    ])
    add_standard_slide("11.3.1 规划成本管理：输入", [
        "1. 项目章程：规定了预先批准的财务资源，确定详细的项目成本。",
        "2. 项目管理计划：",
        "   • 进度管理计划：提供影响成本估算的准则。",
        "   • 风险管理计划：提供识别和监督风险的方法。",
        "3. 事业环境因素：组织文化、市场条件、汇率等。",
        "4. 组织过程资产：财务控制程序、历史信息知识库等。"
    ])
    add_standard_slide("11.3.2 规划成本管理：工具与技术", [
        "1. 专家判断：以往类似项目、行业/应用领域信息、挣值管理等。",
        "2. 数据分析（备选方案分析）：例如自筹、股权或借贷等筹资方案。",
        "3. 会议：项目团队举行规划会议。"
    ])
    add_standard_slide("11.3.3 规划成本管理：输出", [
        "1. 成本管理计划：描述如何规划、安排和控制项目成本。内容包含：",
        "   • 计量单位 (Unit of Measure)",
        "   • 精确度 (Level of Precision)",
        "   • 准确度 (Level of Accuracy)",
        "   • 组织程序链接 (WBS)",
        "   • 控制临界值 (Control Thresholds)",
        "   • 绩效测量规则 (EVM规则)",
        "   • 报告格式"
    ])

    # --- 11.4 估算成本 ---
    add_standard_slide("11.4 估算成本", [
        "估算成本是对完成项目工作所需资源成本进行近似估算的过理。",
        "主要作用：确定项目所需的资金。",
        "特点：随着项目进展，准确性逐步提高（粗略量级估算 -25%~+75%，确定性估算 -5%~+10%）。"
    ])
    add_standard_slide("11.4.1 估算成本：输入", [
        "1. 项目管理计划：成本管理计划、质量管理计划、范围基准。",
        "2. 项目文件：风险登记册、经验教训登记册、资源需求、项目进度计划。",
        "3. 事业环境因素：市场条件、发布商业信息、汇率、地区生产率差异等。",
        "4. 组织过程资产：成本估算政策、模版、历史信息等。"
    ])
    add_standard_slide("11.4.2 估算成本：工具与技术 (1)", [
        "1. 专家判断：以往类似项目信息、估算方法等。",
        "2. 类比估算：基于以往类似项目的参数或属性。",
        "3. 参数估算：利用历史数据之间的统计关系进行计算。",
        "4. 自下而上估算：从下到上逐层汇总工作组成部分的估算。",
        "5. 三点估算：考虑不确定性和风险（三角分布、贝塔分布）。"
    ])
    add_standard_slide("11.4.2 估算成本：工具与技术 (2)", [
        "6. 数据分析：",
        "   • 备选方案分析：比较不同成本方案。",
        "   • 储备分析：包含应急储备，应对“已知-未知”风险。",
        "   • 质量成本：关于质量成本的各种假设。",
        "7. 项目管理信息系统：电子表单、模拟软件等辅助估算。",
        "8. 决策：如投票（增加团队参与度和准确性）。"
    ])
    add_standard_slide("11.4.3 估算成本：输出", [
        "1. 成本估算：覆盖项目全部资源成本，包括直接人工、材料、设备、服务等。",
        "2. 估算依据：说明估算是如何得出的，包括假设条件、已知制约因素、置信水平等。",
        "3. 项目文件更新：假设日志、经验教训登记册、风险登记册等。"
    ])

    # --- 11.5 制定预算 ---
    add_standard_slide("11.5 制定预算", [
        "制定预算是汇总所有单个活动或工作包的估算成本，建立一个经批准的成本基准的过程。",
        "主要作用：确定成本基准，用以监控和控制项目绩效。",
        "项目预算 = 成本基准 + 管理储备。"
    ])
    add_standard_slide("11.5.1 制定预算：输入", [
        "1. 项目管理计划：成本管理计划、资源管理计划、范围基准。",
        "2. 商业文件：可行性研究报告、项目评估报告。",
        "3. 项目文件：估算依据、成本估算、项目进度计划、风险登记册。",
        "4. 协议：采购产品、服务或成果的成本。",
        "5. 事业环境、组织过程资产。"
    ])
    add_standard_slide("11.5.2 制定预算：工具与技术", [
        "1. 专家判断：以往类似项目、财务原则、资金需求等。",
        "2. 成本汇总：先汇总到工作包，再到控制账户，最终得总成本。",
        "3. 数据分析（储备分析）：建立管理储备，应对“未知-未知”风险。",
        "4. 历史信息审核：审核参数估算或类比估算的可靠性。",
        "5. 资金限额平衡：平衡资金限制与计划支出，可调整进度计划。",
        "6. 融资：获取外部资金。"
    ])
    add_standard_slide("11.5.3 制定预算：输出", [
        "1. 成本基准：经批准、按时间段分配的项目预算，不含管理储备。",
        "   • S曲线：按时间分配成本基准得到。",
        "2. 项目资金需求：包括预计支出和预计债务。",
        "3. 项目文件更新：成本估算、项目进度计划、风险登记册。"
    ])

    # --- 11.6 控制成本 ---
    add_standard_slide("11.6 控制成本", [
        "控制成本是监督项目状态，以更新项目成本和管理成本基准变更的过程。",
        "主要作用：在整个项目期间保持对成本基准的维护。",
        "控制重点：分析项目资金支出与已完成工作之间的关系。"
    ])
    add_standard_slide("11.6.1 控制成本：输入", [
        "1. 项目管理计划：成本管理计划、成本基准、绩效测量基准。",
        "2. 项目资金需求：预计支出及债务。",
        "3. 工作绩效数据：发生、支付、开具发票的成本等。",
        "4. 项目文件：经验教训登记册。",
        "5. 组织过程资产。"
    ])
    add_standard_slide("11.6.2 控制成本：工具与技术 (1) - EVM", [
        "1. 专家判断：偏差分析、挣值分析、预测等。",
        "2. 挣值分析 (EVA)：比较计划、挣值与实际。",
        "   • 计划价值 (PV)：计划工作的预算。",
        "   • 挣值 (EV)：已完成工作的预算。",
        "   • 实际成本 (AC)：已完成工作的实际开支。"
    ])
    add_standard_slide("11.6.2 控制成本：工具与技术 (2) - 评价指标", [
        "   • 成本偏差 (CV) = EV - AC (负值表示超支)",
        "   • 进度偏差 (SV) = EV - PV (负值表示落后)",
        "   • 成本绩效指数 (CPI) = EV / AC (小于1表示超支)",
        "   • 进度绩效指数 (SPI) = EV / PV (小于1表示落后)"
    ])
    add_standard_slide("11.6.2 控制成本：工具与技术 (3) - 趋势与预测", [
        "   • 完工估算 (EAC)：基于项目当前绩效预测的总成本。",
        "   • 完工尚需指数 (TCPI)：为实现特定目标需达到的成本绩效。",
        "   • 数据分析：偏差分析、趋势分析、储备分析。"
    ])
    add_standard_slide("11.6.3 控制成本：输出", [
        "1. 工作绩效信息：CV, SV, CPI, SPI, TCPI 等计算结果。",
        "2. 成本预测：预测的 EAC。",
        "3. 变更请求：需采取纠正或预防措施时提出。",
        "4. 项目管理计划及文件更新。"
    ])

    prs.save(output_path)
    print(f"PPT saved to {output_path}")

if __name__ == "__main__":
    out_file = r"e:\教学管理\第11章 项目成本管理.pptx"
    tmpl_file = r"e:\教学管理\第4章 信息系统项目的进度管理.pptx"
    create_ch11_ppt(out_file, tmpl_file)
