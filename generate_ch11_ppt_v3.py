from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_detailed_ch11_ppt(output_path, template_path, img_dir):
    if os.path.exists(template_path):
        prs = Presentation(template_path)
        for _ in range(len(prs.slides)):
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
    else:
        prs = Presentation()

    def add_slide_with_image(title_text, img_name, content_text=None):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title_text
        
        # Add Header & Line (as before)
        header_text = "十一五规划教材                                          《信息系统项目管理》讲义"
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = header_text
        p.font.size = Pt(12)
        slide.shapes.add_connector(1, Inches(0.5), Inches(0.7), Inches(9.5), Inches(0.7))

        if content_text:
            body_shape = slide.placeholders[1]
            body_shape.text = content_text
        
        img_path = os.path.join(img_dir, img_name)
        if os.path.exists(img_path):
            # Centered at the bottom/middle
            # Usually diagrams from PDF are horizontal, so we place them below title
            slide.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(8))

    def add_text_slide(title_text, paragraphs):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title_text
        
        header_text = "十一五规划教材                                          《信息系统项目管理》讲义"
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = header_text
        p.font.size = Pt(12)
        slide.shapes.add_connector(1, Inches(0.5), Inches(0.7), Inches(9.5), Inches(0.7))

        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        for i, text in enumerate(paragraphs):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = text
            p.level = 0
            # If line is a sub-point (starts with bullet), set level 1
            if text.startswith("  ") or text.startswith("- "):
                p.level = 1
                p.text = text.lstrip(" -")

    # 1. Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tx1 = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
    tf1 = tx1.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "信息系统项目管理"
    p1.font.bold = True
    p1.font.size = Pt(44)
    p1.alignment = PP_ALIGN.CENTER
    tx2 = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(2))
    tf2 = tx2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "第11章 项目成本管理 (详细版)"
    p2.font.size = Pt(32)
    p2.alignment = PP_ALIGN.CENTER

    # 2. Intro
    add_text_slide("11.1 项目成本管理概述", [
        "通俗定义：项目成本管理就是决定项目怎么花钱、怎么省钱，并确保在老板批的预算内把活干完。",
        "核心工作：",
        "  • 规划：定规矩（怎么算账）。",
        "  • 估算：报数额（每项活花多少）。",
        "  • 预算：总盘子（一共花多少）。",
        "  • 控制：盯着花（有没有超支）。"
    ])

    # 3. 11.3 Plan Cost Management
    add_text_slide("11.3 规划成本管理 - 定义与作用", [
        "定义：决定如何给项目定预算、怎么花钱、怎么省钱的过程。",
        "作用：就像写一本“财务手册”，规定后续所有成本活动的标准。",
        "常见问题：如果没有这一步，后期报账和核算会非常混乱。"
    ])
    add_slide_with_image("11.3 规划成本管理 - 数据流向", "fig_11_1.png")
    add_text_slide("11.3.1 输入 (老板手里有什么？)", [
        "1. 项目章程：老板给的最初授权和大致预算范围。",
        "2. 项目管理计划：",
        "  • 进度计划：干得快慢直接影响花钱多少。",
        "  • 风险计划：预留多少钱防备意外。",
        "3. 事业环境因素：比如现在的市场行情、外汇汇率等。",
        "4. 组织过程资产：公司以前做类似项目的账目模版。"
    ])
    add_text_slide("11.3.2 工具与技术 (怎么定计划？)", [
        "1. 专家判断：请老财务、资深项目经理凭经验把关。",
        "2. 数据分析（备选方案分析）：",
        "  • 到底是自己做省钱，还是花钱买一套软件更划算？",
        "  • 资金是靠公司自筹，还是去银行贷款更好？",
        "3. 会议：财务、技术、管理人员坐在一起碰头。"
    ])
    add_text_slide("11.3.3 输出 - 成本管理计划", [
        "输出内容包括：",
        "  • 计量单位：是用“人工时”算，还是用“元”算？",
        "  • 精确度：是一分钱都不能差，还是保留到“百元”级别？",
        "  • 控制临界值：超支多少百分比必须报警？(如±5%)",
        "  • 绩效测量规则：确定用什么规则计算挣值 (EVM)。"
    ])

    # 4. 11.4 Estimate Costs
    add_text_slide("11.4 估算成本 - 让每一个环节都有数", [
        "定义：计算完成每一项具体任务需要花多少钱的“概算”。",
        "核心特点：它是一个“进化的过程”。",
        "  • 早期：粗略量级估算 (-25% 到 +75%)。",
        "  • 后期：确定性估算 (-5% 到 +10%)。"
    ])
    add_slide_with_image("11.4 估算成本 - 数据流向", "fig_11_2.png")
    add_text_slide("11.4.2 估算方法 (如何算得准？)", [
        "1. 类比估算：看以前的项目，简单快速，但不一定准。",
        "2. 参数估算：按单价算（如 200元/行代码），适合重复性工作。",
        "3. 三点估算：(最乐观 + 4*最可能 + 最悲观) / 6。考虑了风险。",
        "4. 自下而上估算：把每个零件的钱加起来。最准，但也最慢。",
        "5. 储备分析：预留应急储备金（应对已知的风险）。"
    ])

    # 5. 11.5 Determine Budget
    add_text_slide("11.5 制定预算 - 建立成本基准", [
        "定义：把所有个别活儿的钱加起来，变成一笔正式的“基准拨款”。",
        "核心重点：成本基准 (Cost Baseline)。",
        "这一步会产出项目经理手中的“指挥棒”，衡量是否有钱可花的依据。"
    ])
    add_slide_with_image("11.5 制定预算 - 数据流向", "fig_11_3.png")
    add_slide_with_image("11.5 预算的组成 (重要概念)", "fig_11_4.png", 
                   "注意：项目预算 = 成本基准 + 管理储备。管理储备受老板直接控制，不属于项目经理的日常支配范围。")
    add_text_slide("11.5.2 预算平衡技术", [
        "1. 成本汇总：从下到上凑整。",
        "2. 资金限额平衡：如果公司这个月只有10万预算，项目经理就不能安排需要花20万的活，得调整进度。",
        "3. 历史审核：检查算出来的数跟以前同类项目比，是否靠谱。"
    ])

    # 6. 11.6 Control Costs
    add_text_slide("11.6 控制成本 - 盯着钱是怎么花的", [
        "定义：监督项目干活的情况，看预算有没有超支，进度有没有落后。",
        "关键：不要只看花掉多少钱，要看花这些钱换回了多少活 (EV)。"
    ])
    add_slide_with_image("11.6 控制成本 - 数据流向", "fig_11_6.png")

    # 7. EVM Deep Dive
    add_text_slide("控制成本：挣值分析 (EVM) 详解 (1)", [
        "它是衡量项目健康度的“体检表”，包含三个核心数值：",
        "1. PV (Planned Value) 计划价值：到今天，原计划应该干多少钱的活？",
        "2. EV (Earned Value) 挣值：到今天，看实际干完的活，按原计划值多少钱？",
        "3. AC (Actual Cost) 实际成本：到今天，实际上已经掏了多少钱？"
    ])
    add_slide_with_image("挣值管理曲线 (S-Curve)", "fig_11_7.png")
    add_text_slide("控制成本：挣值分析 (EVM) 详解 (2) - 评价指标", [
        "偏差分析 (加减法)：",
        "  • CV (成本偏差) = EV - AC。正数=省钱，负数=超支。",
        "  • SV (进度偏差) = EV - PV。正数=快，负数=落后。",
        "效率分析 (乘除法)：",
        "  • CPI (成本指数) = EV / AC。>1说明每一块钱产生的活超过一元。",
        "  • SPI (进度指数) = EV / PV。>1说明干得比计划快。"
    ])
    
    # 8. EVM Example
    add_text_slide("【实战演练】EVM 指标计算案例", [
        "项目情况：总预算 10 万元，工期 10 个月。每月计划完成 1 万元的任务。",
        "目前状态：第 4 个月底。",
        "系统数据显示：",
        "  • PV = 4 万元 (计划 4 个月完成 40% 的活)。",
        "  • EV = 3 万元 (实际效率低，只干完了预计值 3 万元的活)。",
        "  • AC = 3.5 万元 (实际这些活花掉了 3.5 万元)。"
    ])
    add_text_slide("【案例分析】项目体检报告", [
        "1. 算结果：",
        "  • CV = 3 - 3.5 = -0.5 万 (超支 5000 元)。",
        "  • SV = 3 - 4 = -1.0 万 (进度落后 1 个月)。",
        "  • CPI = 3 / 3.5 = 0.86 (绩效低下)。",
        "2. 结论报告：",
        "  • 当前项目：既超支了，又落后了！",
        "  • 建议：必须马上分析原因（是不是人手不够，或者材料浪费？），采取纠正措施。"
    ])

    prs.save(output_path)
    print(f"Detailed PPT saved to {output_path}")

if __name__ == "__main__":
    out_file = r"e:\教学管理\第11章 项目成本管理_详细版.pptx"
    tmpl_file = r"e:\教学管理\第4章 信息系统项目的进度管理.pptx"
    img_dir = r"e:\教学管理\ch11_images"
    create_detailed_ch11_ppt(out_file, tmpl_file, img_dir)
