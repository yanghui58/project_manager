from pptx import Presentation
import json

def analyze_pptx(file_path):
    try:
        prs = Presentation(file_path)
        slides_data = []
        
        # Analyze slides
        for i, slide in enumerate(prs.slides):
            if i >= 10: break # Only first 10
            slide_info = {
                "slide_index": i + 1,
                "layout": slide.slide_layout.name,
                "elements": []
            }
            for shape in slide.shapes:
                element = {
                    "name": shape.name,
                    "type": str(shape.shape_type),
                    "text": shape.text if hasattr(shape, "text") and shape.has_text_frame else ""
                }
                if shape.has_table:
                    element["table"] = [[cell.text for cell in row.cells] for row in shape.table.rows]
                slide_info["elements"].append(element)
            slides_data.append(slide_info)
        return slides_data
    except Exception as e:
        return {"error": str(e)}

if __name__ == "__main__":
    pptx1 = r"e:\教学管理\第4章 信息系统项目的进度管理.pptx"
    pptx2 = r"e:\教学管理\第4章 信息系统项目的进度管理-2.pptx"
    
    data = {
        "pptx1": analyze_pptx(pptx1),
        "pptx2": analyze_pptx(pptx2)
    }
    
    with open(r"e:\教学管理\pptx_analysis.json", "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    print("Analysis saved to pptx_analysis.json")
