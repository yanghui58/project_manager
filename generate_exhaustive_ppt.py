from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def generate_exhaustive_ppt(output_path, template_path, img_dir, text_file):
    if os.path.exists(template_path):
        prs = Presentation(template_path)
        for _ in range(len(prs.slides)):
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
    else:
        prs = Presentation()

    def add_slide(title, content_list, img=None):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        
        header_text = "十一五规划教材                                          《信息系统项目管理》讲义"
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = header_text
        p.font.size = Pt(12)
        slide.shapes.add_connector(1, Inches(0.5), Inches(0.7), Inches(9.5), Inches(0.7))

        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        # Determine if we need to split content
        current_tf = tf
        for i, text in enumerate(content_list):
            if i == 0:
                p = current_tf.paragraphs[0]
            else:
                p = current_tf.add_paragraph()
            p.text = text
            if text.strip().startswith(('①','②','③','④','⑤','⑥','•','-')):
                p.level = 1
            else:
                p.level = 0
            p.font.size = Pt(16)

        if img:
            img_path = os.path.join(img_dir, img)
            if os.path.exists(img_path):
                # If image exists, we might need a separate slide or place it carefully
                # To be exhaustive, let's put it on its own slide if it's a big diagram
                img_slide = prs.slides.add_slide(prs.slide_layouts[1])
                img_slide.shapes.title.text = title + " (图表)"
                img_slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(9))

    # Parse text file into sections
    sections = {}
    current_section = "Intro"
    sections[current_section] = []
    
    with open(text_file, 'r', encoding='utf-8') as f:
        for line in f:
            line = line.strip()
            if not line or line.startswith("--- PDF Page"): continue
            if line.startswith(("11.3", "11.4", "11.5", "11.6")):
                current_section = line
                sections[current_section] = []
            else:
                sections[current_section].append(line)

    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tx = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = "信息系统项目管理\n第11章 项目成本管理\n【完整教材版】"
    p.font.size = Pt(44)
    p.alignment = PP_ALIGN.CENTER

    # Generate slides for each section
    for sec_title, content in sections.items():
        # Split content into chunks of 10-12 lines for readability
        chunk_size = 10
        for i in range(0, len(content), chunk_size):
            chunk = content[i : i + chunk_size]
            sub_title = f"{sec_title} ({i//chunk_size + 1})"
            
            # Map images to sections
            img = None
            if "11.3" in sec_title and i == 0: img = "fig_11_1.png"
            elif "11.4" in sec_title and i == 0: img = "fig_11_2.png"
            elif "11.5" in sec_title and i == 0: img = "fig_11_3.png"
            elif "11.6" in sec_title and i == 0: img = "fig_11_6.png"
            
            add_slide(sub_title, chunk, img=img)
            
            # Add ITTO Tables specifically
            if "11.3" in sec_title and i == 0: add_slide("11.3 ITTO汇总", [], img="itto_11_3_v2.png")
            if "11.4" in sec_title and i == 0: add_slide("11.4 ITTO汇总", [], img="itto_11_4_v2.png")
            if "11.5" in sec_title and i == 0: add_slide("11.5 ITTO汇总", [], img="itto_11_5_v2.png")
            if "11.6" in sec_title and i == 0: add_slide("11.6 ITTO汇总", [], img="itto_11_6_v2.png")

    # Add Budget Composition & EVM Tables specifically
    add_slide("11.5.3 项目预算的组成", [], img="fig_11_4.png")
    add_slide("11.6 挣值分析曲线", [], img="fig_11_7.png")
    add_slide("11.6 EVM指标汇总表 (1)", [], img="evm_table_p1.png")
    add_slide("11.6 EVM指标汇总表 (2)", [], img="evm_table_p2.png")

    prs.save(output_path)
    print(f"Exhaustive PPT saved to {output_path}")

if __name__ == "__main__":
    out_file = r"e:\教学管理\第11章 项目成本管理_完整专业版.pptx"
    tmpl_file = r"e:\教学管理\第4章 信息系统项目的进度管理.pptx"
    img_dir = r"e:\教学管理\ch11_images"
    text_file = r"e:\教学管理\ch11_prof_text_final.txt"
    generate_exhaustive_ppt(out_file, tmpl_file, img_dir, text_file)
