import fitz  # PyMuPDF
from pptx import Presentation
import os

pdf_path = r"e:\教学管理\2025-2026-2\信息系统工程管理\课程概述.pdf"
ppt_path = r"e:\教学管理\2025-2026-2\信息系统工程管理\课程概述.pptx"

def pdf_to_ppt(pdf_file, ppt_file):
    print(f"Opening {pdf_file}...")
    pdf_doc = fitz.open(pdf_file)
    
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    
    for page_num in range(len(pdf_doc)):
        print(f"Processing page {page_num + 1}/{len(pdf_doc)}...")
        page = pdf_doc.load_page(page_num)
        
        mat = fitz.Matrix(2.0, 2.0)
        pix = page.get_pixmap(matrix=mat)
        
        img_path = f"temp_page_{page_num}.png"
        pix.save(img_path)
        
        slide = prs.slides.add_slide(blank_slide_layout)
        
        if page_num == 0:
            prs.slide_width = int(page.rect.width * 12700)
            prs.slide_height = int(page.rect.height * 12700)
            
        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        os.remove(img_path)
    
    print(f"Saving presentation to {ppt_file}...")
    prs.save(ppt_file)
    print("Done!")

if __name__ == "__main__":
    pdf_to_ppt(pdf_path, ppt_path)
